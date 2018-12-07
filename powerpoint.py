from pptx import Presentation
import json

PLACEHOLDER_TITLE = 1
PLACEHOLDER_SLIDE_NUMBER = 13
PLACEHOLDER_BODY = 2
PLACEHOLDER_CENTER_TITLE = 3
PLACEHOLDER_SUBTITLE = 4


def parse_ppt(path):
    prs = Presentation(path)
    parse_content = []
    for p, slide in enumerate(prs.slides):
        page = {}
        page["page"] = p
        text = []
        page["text"] = text
        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue

            if shape.placeholder_format.type == PLACEHOLDER_TITLE:
                if not shape.text or not shape.text.strip():
                    continue
                sub_text = {}
                sub_text["type"] = "title"
                sub_text["text"] = shape.text.strip()
                text.append(sub_text)

            if shape.placeholder_format.type == PLACEHOLDER_BODY:
                for paragraph in shape.text_frame.paragraphs:
                    if not paragraph.text or not paragraph.text.strip():
                        continue
                    sub_text = {}
                    sub_text["type"] = "body"
                    sub_text["level"] = paragraph.level
                    sub_text["text"] = paragraph.text.strip()
                    text.append(sub_text)

        parse_content.append(page)

    return parse_content


if __name__ == "__main__":
    path = r"C:\Users\Wii\Google Drive\My Doc\usc\_CS590\8. ER model.pptx"
    with open("parse_ppt.json", "w") as f:
        json.dump(parse_ppt(path), f)